"""
зӯ”иҫ©PPTз”ҹжҲҗжЁЎеқ—
е®Ңж•ҙжЎҶжһ¶пјҡе°ҒйқўвҶ’зӣ®еҪ•вҶ’иғҢжҷҜвҶ’жҠҖжңҜвҶ’дә§е“ҒвҶ’е•ҶдёҡжЁЎејҸвҶ’еӣўйҳҹвҶ’и§„еҲ’вҶ’иҮҙи°ў
"""

import re, os, io
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PPTGenerator:
    """дё“дёҡзӯ”иҫ©PPTз”ҹжҲҗеҷЁ"""

    def __init__(self):
        self.W = 13.333  # 16:9
        self.H = 7.5
        self.slide_num = 0

    def generate_ppt(self, full_text: str, project_name: str, output_path: Path) -> Optional[Path]:
        if not HAS_PPTX:
            return None
        prs = Presentation()
        prs.slide_width = Inches(self.W)
        prs.slide_height = Inches(self.H)

        chapters = self._parse(full_text)
        nums = re.findall(r'(\d+\.?\d*\s*[%пј…дәҝдёҮеҚғ]?\s*(?:е…ғ|зҫҺе…ғ|йЎ№|зҜҮ|еҖҚ))', full_text)
        innovations = re.findall(r'(?:еҲӣж–°зӮ№|ж ёеҝғеҲӣж–°)[пјҡ:\s]*([^\n]{10,60})', full_text)

        # 1. е°Ғйқў
        self._cover(prs, project_name)
        # 2. зӣ®еҪ•
        toc = ["йЎ№зӣ®иғҢжҷҜдёҺз—ӣзӮ№", "ж ёеҝғжҠҖжңҜеҺҹзҗҶдёҺеҲӣж–°", "дә§е“Ғи®ҫи®ЎдёҺеә”з”ЁеңәжҷҜ",
               "еёӮеңәеҲҶжһҗдёҺе•ҶдёҡжЁЎејҸ", "еӣўйҳҹд»Ӣз»ҚдёҺж ёеҝғдјҳеҠҝ", "еҸ‘еұ•и§„еҲ’дёҺиһҚиө„йңҖжұӮ"]
        self._toc(prs, toc)
        # 3. йЎ№зӣ®иғҢжҷҜ (1-2p)
        self._content_slide(prs, "йЎ№зӣ®иғҢжҷҜдёҺиЎҢдёҡз—ӣзӮ№",
            self._extract_bullets(chapters, ["иғҢжҷҜ", "иЎҢдёҡ"]),
            "еёӮеңәз—ӣзӮ№гҖҒж”ҝзӯ–жңәйҒҮгҖҒиЎҢдёҡи¶ӢеҠҝ", nums[:2])
        # 4. ж ёеҝғжҠҖжңҜ (2-3p)
        self._content_slide(prs, "ж ёеҝғжҠҖжңҜеҺҹзҗҶдёҺеҲӣж–°",
            self._extract_bullets(chapters, ["жҠҖжңҜ", "еҲӣж–°", "еҺҹзҗҶ"]),
            "жҠҖжңҜи·ҜзәҝгҖҒеҲӣж–°зӘҒз ҙгҖҒз«һе“ҒеҜ№жҜ”", nums[2:5] if len(nums) > 2 else [])
        # 5. дә§е“Ғ/и§ЈеҶіж–№жЎҲ
        prod = self._extract_bullets(chapters, ["дә§е“Ғ", "еә”з”Ё", "еңәжҷҜ", "и®ҫи®Ў"])
        if prod:
            self._content_slide(prs, "дә§е“Ғи®ҫи®ЎдёҺеә”з”ЁеңәжҷҜ", prod, "дә§е“ҒеҪўжҖҒгҖҒеә”з”ЁиҗҪең°", [])
        # 6. е•ҶдёҡжЁЎејҸ
        mkt = self._extract_bullets(chapters, ["еёӮеңә", "е•Ҷдёҡ", "жЁЎејҸ"])
        if mkt:
            self._content_slide(prs, "еёӮеңәеҲҶжһҗдёҺе•ҶдёҡжЁЎејҸ", mkt, "е•ҶдёҡжЁЎејҸгҖҒзӣҲеҲ©и·Ҝеҫ„гҖҒеёӮеңәз©әй—ҙ", nums[4:6] if len(nums) > 4 else [])
        # 7. еӣўйҳҹ
        team = self._extract_bullets(chapters, ["еӣўйҳҹ", "жҲҗе‘ҳ", "д»Ӣз»Қ"])
        if team:
            self._team_slide(prs, team)
        # 8. еҸ‘еұ•и§„еҲ’
        plan = self._extract_bullets(chapters, ["и§„еҲ’", "еҸ‘еұ•", "жңӘжқҘ", "иҙўеҠЎ", "иһҚиө„"])
        if plan:
            self._content_slide(prs, "еҸ‘еұ•и§„еҲ’дёҺиһҚиө„йңҖжұӮ", plan, "жҠҖжңҜи·ҜзәҝгҖҒеёӮеңәжӢ“еұ•гҖҒиө„йҮ‘з”ЁйҖ”", [])
        # 9. иҮҙи°ў
        self._thanks(prs)

        prs.save(str(output_path))
        return output_path

    def _parse(self, text: str) -> dict:
        chapters = {}; current = None
        for line in text.split('\n'):
            m = re.match(r'^## (.+)', line)
            if m: current = m.group(1); chapters[current] = []
            elif current and line.strip(): chapters[current].append(line.strip())
        return chapters

    def _extract_bullets(self, chapters: dict, keywords: List[str], max_items: int = 6) -> list:
        items = []
        for k, v in chapters.items():
            if any(kw in k for kw in keywords):
                for line in v:
                    # еҪ»еә•иҝҮж»Өж Үи®°иЎҢ
                    if 'гҖҗеҫ…иЎҘе……гҖ‘' in line or 'гҖҗжқҘжәҗгҖ‘' in line:
                        continue
                    clean = re.sub(r'^[-*#>\d\.\sгҖҗгҖ‘]+', '', line).strip()
                    # иҝҮж»ӨеӨӘзҹӯжҲ–зәҜж Үи®°
                    if len(clean) < 10 or clean.startswith('гҖҗ') or clean.startswith('>'):
                        continue
                    items.append(clean[:130])
        # еҺ»йҮҚ
        seen = set(); result = []
        for item in items:
            key = item[:25]
            if key not in seen: seen.add(key); result.append(item)
        # иҝ”еӣһжӣҙе®Ңж•ҙзҡ„еҲ—иЎЁ
        return result[:max_items]

    def _add_slide_number(self, slide):
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        self.slide_num += 1
        tb = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(1), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.text = str(self.slide_num); p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF); p.alignment = PP_ALIGN.RIGHT

    def _cover(self, prs, project_name):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background; fill = bg.fill; fill.solid()
        fill.fore_color.rgb = RGBColor(0x0A, 0x2F, 0x5A)

        # Title
        tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.9), Inches(2))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = project_name; p.font.size = Pt(40)
        p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); p.alignment = PP_ALIGN.CENTER

        # Subtitle
        tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(10.9), Inches(0.6))
        p2 = tb2.text_frame.paragraphs[0]; p2.text = "з«һиөӣзӯ”иҫ©жұҮжҠҘ"
        p2.font.size = Pt(22); p2.font.color.rgb = RGBColor(0xFF, 0x6B, 0x35); p2.alignment = PP_ALIGN.CENTER

        # Accent line
        line = slide.shapes.add_shape(1, Inches(5), Inches(5.3), Inches(3.333), Inches(0.04))
        line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xFF, 0x6B, 0x35); line.line.fill.background()

    def _toc(self, prs, items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Title bar
        self._title_bar(slide, "жұҮжҠҘзӣ®еҪ•")
        # Numbered items
        tb = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(5))
        tf = tb.text_frame; tf.word_wrap = True
        for i, item in enumerate(items, 1):
            p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
            p.text = f"{i:02d}   {item}"; p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x1F, 0x29, 0x37); p.space_after = Pt(14)
        self._add_slide_number(slide)

    def _title_bar(self, slide, title):
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(self.W), Inches(1.25))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0x0A, 0x2F, 0x5A); bar.line.fill.background()
        tf = bar.text_frame; tf.word_wrap = True; tf.margin_left = Inches(1)
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(26)
        p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    def _content_slide(self, prs, title, bullets, subtitle="", highlight_nums=None):
        # и·іиҝҮж— еҶ…е®№зҡ„йЎө
        if not bullets:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._title_bar(slide, title)

        y = 1.8
        # Subtitle
        if subtitle:
            tb = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.4))
            p = tb.text_frame.paragraphs[0]; p.text = subtitle
            p.font.size = Pt(13); p.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
            y += 0.6

        # Bullet points - жӣҙеӨҡгҖҒжӣҙзҙ§еҮ‘
        if bullets:
            tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11.3), Inches(5.5 - (y - 1.8)))
            tf2 = tb2.text_frame; tf2.word_wrap = True
            for i, b in enumerate(bullets[:8]):
                p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
                p.text = f"в–ё {b[:140]}"; p.font.size = Pt(15)
                p.font.color.rgb = RGBColor(0x1F, 0x29, 0x37); p.space_after = Pt(10)

        # Highlight numbers at the bottom
        if highlight_nums:
            y_bot = 6.3
            for i, num in enumerate(highlight_nums[:4]):
                x = 1.2 + i * 3
                card = slide.shapes.add_shape(1, Inches(x), Inches(y_bot), Inches(2.7), Inches(0.9))
                card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
                card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
                tb3 = slide.shapes.add_textbox(Inches(x+0.2), Inches(y_bot+0.15), Inches(2.3), Inches(0.6))
                p3 = tb3.text_frame.paragraphs[0]; p3.text = num
                p3.font.size = Pt(18); p3.font.bold = True
                p3.font.color.rgb = RGBColor(0xFF, 0x6B, 0x35); p3.alignment = PP_ALIGN.CENTER

        # Speaker notes
        notes_text = f"гҖҗжј”и®ІжҸҗзӨәгҖ‘\n- {title}пјҡе…ұ{len(bullets)}дёӘиҰҒзӮ№\n- е»әи®®ж—¶й•ҝпјҡ1.5-2еҲҶй’ҹ\n- йҮҚзӮ№ејәи°ғж•°жҚ®еҜ№жҜ”е’Ңе·®ејӮеҢ–дјҳеҠҝ"
        slide.notes_slide.notes_text_frame.text = notes_text
        self._add_slide_number(slide)

    def _team_slide(self, prs, team_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._title_bar(slide, "еӣўйҳҹд»Ӣз»ҚдёҺж ёеҝғдјҳеҠҝ")

        members = []
        for item in team_items:
            if any(kw in item for kw in ['еҚҡеЈ«', 'зЎ•еЈ«', 'ж•ҷжҺҲ', 'еҜјеёҲ', 'иҙҹиҙЈдәә', 'жҢҮеҜј', 'жҲҗе‘ҳ']):
                # жҢүйҖ—еҸ·жӢҶжҲҗеӨҡиЎҢжҳҫзӨә
                parts = [p.strip() for p in item.split(',') if p.strip()]
                members.append(parts[:4])  # жңҖеӨҡ4иЎҢ
        members = members[:6]

        for i, parts in enumerate(members):
            col, row = i % 3, i // 3
            x, y = 1.2 + col * 3.8, 2.2 + row * 2.5
            card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.4), Inches(2))
            card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
            accent = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.4), Inches(0.06))
            accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor(0xFF, 0x6B, 0x35)
            accent.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.3), Inches(3), Inches(1.5))
            tf = tb.text_frame; tf.word_wrap = True
            for j, part in enumerate(parts):
                p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                p.text = part[:40]
                p.font.size = Pt(10) if j > 0 else Pt(11)
                p.font.bold = (j == 0)
                p.font.color.rgb = RGBColor(0x1F, 0x29, 0x37) if j == 0 else RGBColor(0x6B, 0x72, 0x80)
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(2)

        slide.notes_slide.notes_text_frame.text = "гҖҗжј”и®ІжҸҗзӨәгҖ‘\n- ејәи°ғеӣўйҳҹдә’иЎҘжҖ§е’ҢжҠҖжңҜз§ҜзҙҜ\n- зӘҒеҮәжҢҮеҜјиҖҒеёҲзҡ„иЎҢдёҡеҪұе“ҚеҠӣ"
        self._add_slide_number(slide)

    def _thanks(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background; fill = bg.fill; fill.solid()
        fill.fore_color.rgb = RGBColor(0x0A, 0x2F, 0x5A)
        tb = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
        p = tb.text_frame.paragraphs[0]; p.text = "ж„ҹи°ўиҒҶеҗ¬"
        p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        tb2 = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9), Inches(1))
        p2 = tb2.text_frame.paragraphs[0]; p2.text = "жҒіиҜ·еҗ„дҪҚиҜ„е§”иҖҒеёҲжү№иҜ„жҢҮжӯЈ"
        p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(0xFF, 0x6B, 0x35); p2.alignment = PP_ALIGN.CENTER
        slide.notes_slide.notes_text_frame.text = "гҖҗжј”и®ІжҸҗзӨәгҖ‘\n- жңҖеҗҺж„ҹи°ўиҜ„е§”\n- з•ҷеҮәQ&Aж—¶й—ҙ"
        self._add_slide_number(slide)

    def analyze_competitors(self, full_text: str, market_data: str = "") -> List:
        @dataclass
        class C: name: str; key_params: dict = None
        known = ["иӢұйЈһеҮҢ","Navitas","Nortek","Vertiv","иҘҝй—Ёеӯҗ","ABB","еҚҺдёә","дёӯе…ҙ","Sigma-Aldrich"]
        result = []
        for n in known:
            if n in full_text: result.append(C(name=n))
        return result[:8]

    def format_competitor_table(self, competitors: List) -> str:
        if not competitors: return ""
        lines = ["## з«һе“ҒеҜ№ж ҮеҲҶжһҗ", "", "| з«һе“Ғ | иҜҙжҳҺ |", "|------|------|"]
        for c in competitors: lines.append(f"| {c.name} | дё»иҰҒз«һдәүеҜ№жүӢ |")
        return "\n".join(lines)
